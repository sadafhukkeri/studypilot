"""Text extractors for various file types and YouTube."""
import io
import re
import base64
from typing import Tuple

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi


def extract_pdf(data: bytes) -> str:
    text_parts = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts).strip()


def extract_docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def extract_txt(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1", errors="ignore")


def extract_youtube(url: str) -> Tuple[str, str]:
    """Returns (transcript_text, video_id)."""
    video_id = _parse_youtube_id(url)
    if not video_id:
        raise ValueError("Could not parse YouTube video ID")
    try:
        api = YouTubeTranscriptApi()
        fetched = api.fetch(video_id)
        # fetched is iterable of FetchedTranscriptSnippet objects with .text
        text = " ".join(snippet.text for snippet in fetched)
        return text, video_id
    except Exception as e:
        # Fallback to legacy static method
        try:
            chunks = YouTubeTranscriptApi.get_transcript(video_id)
            return " ".join(c["text"] for c in chunks), video_id
        except Exception:
            raise ValueError(f"Could not fetch transcript: {e}")


def _parse_youtube_id(url: str) -> str:
    patterns = [
        r"(?:v=|\/)([0-9A-Za-z_-]{11})(?:[?&]|$)",
        r"youtu\.be\/([0-9A-Za-z_-]{11})",
        r"^([0-9A-Za-z_-]{11})$",
    ]
    for p in patterns:
        m = re.search(p, url)
        if m:
            return m.group(1)
    return ""


def extract_text_from_upload(filename: str, content: bytes) -> str:
    name = filename.lower()
    if name.endswith(".pdf"):
        return extract_pdf(content)
    if name.endswith(".docx") or name.endswith(".doc"):
        return extract_docx(content)
    if name.endswith(".pptx") or name.endswith(".ppt"):
        return extract_pptx(content)
    if name.endswith(".txt") or name.endswith(".md"):
        return extract_txt(content)
    raise ValueError(f"Unsupported file type for text extraction: {filename}")


def image_to_base64(content: bytes) -> str:
    return base64.b64encode(content).decode("utf-8")
